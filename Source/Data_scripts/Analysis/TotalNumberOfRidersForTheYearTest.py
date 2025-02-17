import pandas as pd
import matplotlib
matplotlib.use('Agg')  # Disable GPU acceleration to prevent BSOD
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
from pathlib import Path

# --- 1. Define Paths ---
base_dir = Path(__file__).resolve().parents[3]
input_dir = base_dir / "Source" / "Data" / "Raw"
output_dir = base_dir / "Source" / "Data" / "reports"
output_dir.mkdir(parents=True, exist_ok=True)

# Define file paths
file_path = input_dir / "MTA_Subway_Hourly_Ridership__2020-2024.csv"
output_file = output_dir / "MTA_Station_Ridership_Analysis.xlsx"
excel_path = output_dir / "Average_Ridership_Data.xlsx"
ppt_path = output_dir / "Average_Ridership_Presentation.pptx"

if not file_path.exists():
    raise FileNotFoundError(f"🚨 File not found: {file_path}")

# --- 2. Load Data Efficiently ---
chunks = pd.read_csv(file_path, chunksize=10000, parse_dates=["transit_timestamp"], infer_datetime_format=True, low_memory=False)

df_list = []
agg_ridership = pd.DataFrame()
agg_stations = pd.DataFrame()

for chunk in chunks:
    chunk.dropna(subset=["transit_timestamp"], inplace=True)
    chunk["year"] = chunk["transit_timestamp"].dt.year
    chunk["hour"] = chunk["transit_timestamp"].dt.hour
    chunk["AM_PM"] = chunk["hour"].apply(lambda x: f"{x % 12 if x % 12 != 0 else 12} {'AM' if x < 12 else 'PM'}")
    
    chunk_agg = chunk.groupby(["year", "AM_PM"])['ridership'].mean().reset_index()
    agg_ridership = pd.concat([agg_ridership, chunk_agg]).groupby(["year", "AM_PM"]).mean().reset_index()
    
    chunk_station_agg = chunk.groupby(["year", "station_complex"])['ridership'].sum().reset_index()
    agg_stations = pd.concat([agg_stations, chunk_station_agg]).groupby(["year", "station_complex"]).sum().reset_index()
    
    df_list.append(chunk)

df = pd.concat(df_list)

total_ridership_per_year = df.groupby("year")["ridership"].sum()
official_ridership_2023 = total_ridership_per_year.get(2023, 0)
official_ridership_2024 = total_ridership_per_year.get(2024, 0)

# --- 3. Compute Percentage Per Station ---
stations_2023 = df[df["year"] == 2023].groupby("station_complex")["ridership"].sum().reset_index()
stations_2024 = df[df["year"] == 2024].groupby("station_complex")["ridership"].sum().reset_index()

stations_2023["percentage"] = (stations_2023["ridership"] / official_ridership_2023) * 100 if official_ridership_2023 > 0 else 0
stations_2024["percentage"] = (stations_2024["ridership"] / official_ridership_2024) * 100 if official_ridership_2024 > 0 else 0

# --- 4. Save to Excel ---
with pd.ExcelWriter(output_file, engine="xlsxwriter") as writer:
    stations_2023.to_excel(writer, sheet_name="2023 Ridership", index=False)
    stations_2024.to_excel(writer, sheet_name="2024 Ridership", index=False)
    agg_ridership.to_excel(writer, sheet_name="Avg Ridership", index=False)
    agg_stations.to_excel(writer, sheet_name="All Stations Ridership", index=False)

# --- 5. Create and Save Charts ---
fig, ax = plt.subplots(figsize=(10, 5))
ax.barh(stations_2023["station_complex"], stations_2023["ridership"], label="2023", color="blue")
ax.barh(stations_2024["station_complex"], stations_2024["ridership"], label="2024", color="red", alpha=0.7)
ax.set_xlabel("Ridership")
ax.set_ylabel("Station Complex")
ax.set_title("Top 10 Subway Stations Ridership (2023 vs 2024)")
ax.legend()
plt.gca().invert_yaxis()
chart_path = output_dir / "top_10_ridership_chart.png"
plt.savefig(chart_path, bbox_inches="tight")
plt.close(fig)

# --- 6. Create PowerPoint ---
ppt = Presentation()
for year in [2023, 2024]:
    year_data = agg_ridership[agg_ridership["year"] == year]
    
    plt.figure(figsize=(10, 6))
    plt.plot(year_data["AM_PM"], year_data["ridership"], marker='o', linestyle='-', label=f"{year}")
    plt.title(f"Average Ridership Per Hour - {year}", fontsize=16)
    plt.xlabel("Time (EST)", fontsize=14)
    plt.ylabel("Avg Riders", fontsize=14)
    plt.xticks(rotation=45)
    plt.grid(True, linestyle='--', alpha=0.6)
    plt.legend()
    plt.tight_layout()
    
    chart_filename = output_dir / f"Average_Ridership_Per_Hour_{year}.png"
    plt.savefig(chart_filename)
    plt.close('all')
    
    slide = ppt.slides.add_slide(ppt.slide_layouts[5])
    title = slide.shapes.title
    if title:
        title.text = f"Average Ridership Per Hour - {year}"
    
    left = Inches(1)
    top = Inches(1.5)
    height = Inches(5)
    slide.shapes.add_picture(str(chart_filename), left, top, height=height)

ppt.save(ppt_path)

print(f"✅ Data analysis complete. Excel and PowerPoint reports saved.")
