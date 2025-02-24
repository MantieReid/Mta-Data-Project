import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
import os
from pathlib import Path
from datetime import datetime

def process_year_data(chunks, year):
    """Process data for a specific year"""
    daily_ridership = pd.DataFrame()
    
    for chunk in chunks:
        # Filter data for the specified year and create a copy
        mask = chunk["transit_timestamp"].dt.year == year
        chunk_filtered = chunk[mask].copy()
        
        if not chunk_filtered.empty:
            # Add new columns to the copy
            chunk_filtered.loc[:, "date"] = chunk_filtered["transit_timestamp"].dt.date
            chunk_filtered.loc[:, "day_of_week"] = chunk_filtered["transit_timestamp"].dt.day_name()
            
            # Sum ridership by date
            daily_sum = chunk_filtered.groupby(["date", "day_of_week"])["ridership"].sum().reset_index()
            daily_ridership = pd.concat([daily_ridership, daily_sum])

    # Group by date and sum to get true daily totals
    daily_ridership = daily_ridership.groupby(["date", "day_of_week"])["ridership"].sum().reset_index()

    # Calculate the average ridership for each day of the week
    avg_ridership = daily_ridership.groupby("day_of_week")["ridership"].mean().reindex([
        "Monday", "Tuesday", "Wednesday", "Thursday", "Friday", "Saturday", "Sunday"
    ])
    
    return avg_ridership

def create_chart(avg_ridership, year, base_dir):
    """Create and save a bar chart for the specified year"""
    plt.figure(figsize=(12, 6))
    bars = plt.bar(avg_ridership.index, avg_ridership.values, color='#1f77b4', alpha=0.8)
    plt.xlabel("Day of the Week", fontsize=10, fontweight='bold')
    plt.ylabel("Average Ridership", fontsize=10, fontweight='bold')
    plt.title(f"Average Daily Subway Ridership by Day of Week ({year})", 
              fontsize=12, fontweight='bold', pad=20)
    plt.xticks(rotation=45, ha="right")
    plt.grid(axis="y", linestyle="--", alpha=0.7)

    # Add value labels on top of bars
    def format_with_commas(x):
        return f'{x:,.0f}'

    for bar in bars:
        height = bar.get_height()
        plt.text(bar.get_x() + bar.get_width()/2., height,
                format_with_commas(height),
                ha='center', va='bottom')

    plt.tight_layout()

    # Save the chart
    chart_path = os.path.join(base_dir, "Data", "processed", f"ridership_chart_{year}.png")
    plt.savefig(chart_path, dpi=300, bbox_inches='tight')
    plt.close()
    
    return chart_path

def save_to_excel(avg_ridership_2023, avg_ridership_2024, base_dir):
    """Save both years' data to Excel"""
    excel_path = os.path.join(base_dir, "Data", "processed", "MTA_Subway_Ridership_Weekday_Stats_average.xlsx")
    
    with pd.ExcelWriter(excel_path, engine='xlsxwriter') as writer:
        workbook = writer.book
        number_format = workbook.add_format({'num_format': '#,##0'})

        # Save 2023 data
        avg_ridership_df_2023 = avg_ridership_2023.reset_index()
        avg_ridership_df_2023.columns = ["Day of the Week", "Average Ridership"]
        avg_ridership_df_2023.to_excel(writer, index=False, sheet_name='2023 Average Ridership')
        
        # Format 2023 sheet
        worksheet_2023 = writer.sheets['2023 Average Ridership']
        worksheet_2023.set_column('B:B', 18, number_format)
        worksheet_2023.set_column('A:A', 15)
        worksheet_2023.add_table(0, 0, len(avg_ridership_df_2023), 1, {
            'columns': [
                {'header': 'Day of the Week'},
                {'header': 'Average Ridership', 'format': number_format}
            ],
            'style': 'Table Style Medium 2'
        })

        # Save 2024 data
        avg_ridership_df_2024 = avg_ridership_2024.reset_index()
        avg_ridership_df_2024.columns = ["Day of the Week", "Average Ridership"]
        avg_ridership_df_2024.to_excel(writer, index=False, sheet_name='2024 Average Ridership')
        
        # Format 2024 sheet
        worksheet_2024 = writer.sheets['2024 Average Ridership']
        worksheet_2024.set_column('B:B', 18, number_format)
        worksheet_2024.set_column('A:A', 15)
        worksheet_2024.add_table(0, 0, len(avg_ridership_df_2024), 1, {
            'columns': [
                {'header': 'Day of the Week'},
                {'header': 'Average Ridership', 'format': number_format}
            ],
            'style': 'Table Style Medium 2'
        })
    
    return excel_path

def create_powerpoint(chart_path_2023, chart_path_2024, base_dir):
    """Create PowerPoint presentation with both years' charts"""
    current_time = datetime.now()
    date_time_str = current_time.strftime("%B %d, %Y %I-%M %p")

    # Create PowerPoint presentation with date in filename
    base_filename = f"MTA_Subway_Ridership_Weekday_Stats_Average{date_time_str}.pptx"
    ppt_dir = os.path.join(base_dir, "Data", "charts", "Powerpoint_format")
    os.makedirs(ppt_dir, exist_ok=True)
    ppt_path = os.path.join(ppt_dir, base_filename)

    # Handle duplicate files
    counter = 1
    while os.path.exists(ppt_path):
        base_filename = f"average_ridership_For_All_Days_analysis_{date_time_str}_{counter}.pptx"
        ppt_path = os.path.join(ppt_dir, base_filename)
        counter += 1

    # Create the presentation
    prs = Presentation()
    
    # Add 2023 slide
    slide_layout = prs.slide_layouts[5]
    slide_2023 = prs.slides.add_slide(slide_layout)
    title_2023 = slide_2023.shapes.title
    title_2023.text = "Average Daily Subway Ridership by Day of Week (2023)"
    slide_2023.shapes.add_picture(chart_path_2023, Inches(1), Inches(2.5), width=Inches(8))

    # Add 2024 slide
    slide_2024 = prs.slides.add_slide(slide_layout)
    title_2024 = slide_2024.shapes.title
    title_2024.text = "Average Daily Subway Ridership by Day of Week (2024)"
    slide_2024.shapes.add_picture(chart_path_2024, Inches(1), Inches(2.5), width=Inches(8))

    # Save the PowerPoint file
    prs.save(ppt_path)
    return ppt_path

def main():
    # Set up paths
    base_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", ".."))
    file_path = os.path.join(base_dir, "Data", "Raw", "MTA_Subway_Hourly_Ridership__2020-2024.csv")
    
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"CSV file not found at path: {file_path}")

    # Create output directories
    for dir_path in [
        os.path.join(base_dir, "Data", "reports"),
        os.path.join(base_dir, "Data", "reports",)
    ]:
        os.makedirs(dir_path, exist_ok=True)

    # Load and process data
    date_format = '%m/%d/%Y %I:%M:%S %p'
    chunks = pd.read_csv(file_path, 
                        chunksize=100000,
                        low_memory=False,
                        parse_dates=['transit_timestamp'],
                        date_format=date_format)
    
    # Process both years
    avg_ridership_2023 = process_year_data(chunks, 2023)
    
    # Reset file pointer for 2024
    chunks = pd.read_csv(file_path, 
                        chunksize=100000,
                        low_memory=False,
                        parse_dates=['transit_timestamp'],
                        date_format=date_format)
    avg_ridership_2024 = process_year_data(chunks, 2024)

    # Create charts
    chart_path_2023 = create_chart(avg_ridership_2023, 2023, base_dir)
    chart_path_2024 = create_chart(avg_ridership_2024, 2024, base_dir)

    # Save to Excel
    excel_path = save_to_excel(avg_ridership_2023, avg_ridership_2024, base_dir)

    # Create PowerPoint
    ppt_path = create_powerpoint(chart_path_2023, chart_path_2024, base_dir)

    # Print output paths
    print("✅ Excel file saved at:", excel_path)
    print("✅ PowerPoint file saved at:", ppt_path)

if __name__ == "__main__":
    main()