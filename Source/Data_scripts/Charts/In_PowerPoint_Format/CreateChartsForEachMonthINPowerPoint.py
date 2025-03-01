import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
import os
import gc  # Import garbage collection module
import io  # Import io for BytesIO
import traceback  # For better error reporting

# Define chunk size for processing
CHUNK_SIZE = 10000000  # Adjust based on available RAM

current_dir = os.getcwd()
project_root = os.path.abspath(os.path.join(current_dir, "..", "..", "..", ".."))
file_path = os.path.join(project_root, "Source", "Data", "Raw", "MTA_Subway_Hourly_Ridership__2020-2024.csv")
file_path_OutPut = os.path.join(project_root, "Source", "Data", "reports")

# Add watermark text
WATERMARK_TEXT = "Created By Mantie Reid II"

# Print file path for debugging
print(f"Looking for data file at: {file_path}")
if not os.path.exists(file_path):
    print(f"ERROR: File not found at {file_path}")
    # Try to find the file in the current directory
    alternative_path = "MTA_Subway_Hourly_Ridership__2020-2024.csv"
    if os.path.exists(alternative_path):
        print(f"Found file at alternative location: {alternative_path}")
        file_path = alternative_path
    else:
        print("Please provide the correct path to the CSV file.")
        exit(1)

# Define all possible hours for completeness
hours = [f"{h % 12 if h % 12 != 0 else 12} {'AM' if h < 12 else 'PM'}" for h in range(24)]
all_hours = pd.DataFrame({"AM_PM": hours})
hour_indices = {hour: i for i, hour in enumerate(hours)}  # For sorting

# Function to sanitize station names for filenames
def sanitize_name(name):
    if not isinstance(name, str):
        return "Unknown_Station"
    return name.replace("/", "-").replace("\\", "-").replace("*", "-") \
              .replace("[", "(").replace("]", ")").replace(":", "-") \
              .replace("?", "-").replace("'", "-").replace(",", "-") \
              .replace(".", "-").strip()[:31]

# Create a tracking dictionary for processed stations
processed_stations = {}

# Required columns - only read what we need
required_cols = ['transit_timestamp', 'station_complex_id', 'station_complex', 'ridership']

# Read first few rows to inspect the timestamp format
try:
    sample_df = pd.read_csv(file_path, nrows=5)
    print("Sample data:")
    print(sample_df[['transit_timestamp']].head())
    
    # Try to detect timestamp format
    sample_timestamp = sample_df['transit_timestamp'].iloc[0] if not sample_df.empty else None
    print(f"Sample timestamp: {sample_timestamp}")
except Exception as e:
    print(f"Error reading sample data: {str(e)}")

# Process data for both 2023 and 2024
for year in [2023]:
    print(f"Processing data for year {year}")
    
    # Aggregation structure to store cumulative sums and counts
    month_station_data = {}  # {(month, station_id): {"name": str, "sums": {hour: sum}, "counts": {hour: count}}}
    
    # Track which months we've seen
    months_seen = set()
    
    try:
        # Read and process in chunks to reduce memory usage
        chunks_processed = 0
        
        # Using correct timestamp format and low_memory=False to handle mixed types
        for chunk in pd.read_csv(file_path, chunksize=CHUNK_SIZE, usecols=required_cols, low_memory=False):
            chunks_processed += 1
            
            if chunks_processed % 5 == 0:
                print(f"Processing chunk {chunks_processed}")
            
            # Convert timestamp with the correct format (MM/DD/YYYY HH:MM:SS AM/PM)
            try:
                # Print a sample timestamp from this chunk
                if not chunk.empty:
                    print(f"Sample timestamp from chunk {chunks_processed}: {chunk['transit_timestamp'].iloc[0]}")
                
                chunk["transit_timestamp"] = pd.to_datetime(
                    chunk["transit_timestamp"], 
                    format="%m/%d/%Y %I:%M:%S %p",  # Correct format for "10/18/2022 07:00:00 PM"
                    errors="coerce"
                )
                
                # Print the converted timestamp
                if not chunk.empty:
                    print(f"Converted to: {chunk['transit_timestamp'].iloc[0]}")
            except Exception as e:
                print(f"Error converting timestamps in chunk {chunks_processed}: {e}")
                # Try alternative approach
                chunk["transit_timestamp"] = pd.to_datetime(chunk["transit_timestamp"], errors="coerce")
            
            # Drop rows with invalid timestamps
            valid_before = len(chunk)
            chunk = chunk.dropna(subset=["transit_timestamp"])
            valid_after = len(chunk)
            
            if valid_before > 0:
                print(f"Chunk {chunks_processed}: {valid_after}/{valid_before} valid timestamps ({valid_after/valid_before:.1%})")
            
            if chunk.empty:
                print(f"Chunk {chunks_processed} has no valid timestamps after filtering, skipping")
                continue
                
            # Extract year and filter data
            chunk_year = chunk["transit_timestamp"].dt.year
            
            # Print year distribution to debug
            if chunks_processed <= 2:
                year_counts = chunk_year.value_counts()
                print(f"Years in chunk {chunks_processed}:")
                print(year_counts)
            
            # Filter by year
            year_mask = chunk_year == year
            year_data_count = year_mask.sum()
            
            if year_data_count == 0:
                # Skip chunk if no data for this year
                print(f"Chunk {chunks_processed} has no data for year {year}, skipping")
                continue
            
            print(f"Found {year_data_count} rows for year {year} in chunk {chunks_processed}")
                
            # Work with filtered data - use boolean indexing instead of .copy() to save memory
            filtered_chunk = chunk.loc[year_mask].copy()
            
            # Extract month and hour
            months = filtered_chunk["transit_timestamp"].dt.month
            hours_data = filtered_chunk["transit_timestamp"].dt.hour
            
            # Convert to AM/PM format
            filtered_chunk["AM_PM"] = hours_data.apply(lambda h: f"{h % 12 if h % 12 != 0 else 12} {'AM' if h < 12 else 'PM'}")
            
            # Update months we've seen
            new_months = sorted(months.unique())
            months_seen.update(new_months)
            print(f"Found months in chunk {chunks_processed}: {new_months}")
            
            # Manual approach for aggregation instead of using groupby
            # This avoids the multi-index issues with the agg_result
            try:
                # Group data by station, month, and hour
                for _, row in filtered_chunk.iterrows():
                    station_id = row["station_complex_id"]
                    station_name = row["station_complex"]
                    month = row["transit_timestamp"].month
                    am_pm = row["AM_PM"]
                    ridership = row["ridership"]
                    
                    # Skip NaN ridership values
                    if pd.isna(ridership):
                        continue
                        
                    key = (month, station_id)
                    
                    # Initialize if needed
                    if key not in month_station_data:
                        month_station_data[key] = {
                            "name": station_name,
                            "sums": {h: 0 for h in hours},
                            "counts": {h: 0 for h in hours}
                        }
                    
                    # Update sums and counts
                    month_station_data[key]["sums"][am_pm] += ridership
                    month_station_data[key]["counts"][am_pm] += 1
                
                print(f"Successfully processed {len(filtered_chunk)} rows in chunk {chunks_processed}")
                print(f"Current station-month combinations: {len(month_station_data)}")
                
                # If current dictionary size is large, print status
                if len(month_station_data) > 0 and chunks_processed % 2 == 0:
                    print(f"First 5 months/stations: {list(month_station_data.keys())[:5]}")
                    # Sample entry to debug data structure
                    sample_key = next(iter(month_station_data))
                    print(f"Sample data for {sample_key}:")
                    print(f"  Station name: {month_station_data[sample_key]['name']}")
                    # Show a sample hour with data
                    for hour in hours:
                        if month_station_data[sample_key]['counts'][hour] > 0:
                            print(f"  {hour}: {month_station_data[sample_key]['sums'][hour]} / {month_station_data[sample_key]['counts'][hour]} = {month_station_data[sample_key]['sums'][hour] / month_station_data[sample_key]['counts'][hour]}")
                            break
                
            except Exception as e:
                print(f"Error processing rows in chunk {chunks_processed}: {e}")
                traceback.print_exc()
            
            # Force garbage collection after each chunk
            del chunk, filtered_chunk
            gc.collect()
    
    except Exception as e:
        print(f"Error processing chunks: {str(e)}")
        print("Traceback:")
        traceback.print_exc()
        continue
    
    print(f"Finished processing chunks for year {year}")
    print(f"Months with data: {sorted(months_seen)}")
    print(f"Number of station-month combinations: {len(month_station_data)}")
    
    # Create presentations for each month
    months_with_data = sorted(months_seen)
    
    if not months_with_data:
        print(f"No months found for year {year}, skipping presentation creation")
        continue
        
    print(f"Creating presentations for {len(months_with_data)} months in {year}: {months_with_data}")
    
    for month in months_with_data:
        print(f"Creating presentation for {month}/{year}")
        
        # Get all stations for this month
        month_stations = {station_id: data for (m, station_id), data in month_station_data.items() if m == month}
        
        if not month_stations:
            print(f"No stations have data for {month}/{year}")
            continue
            
        print(f"Found {len(month_stations)} stations with data for {month}/{year}")
        
        # Create PowerPoint presentation
        ppt = Presentation()
        chart_count = 0
        
        # Process each station
        for station_id, station_info in month_stations.items():
            # Skip if already processed
            station_key = f"{station_id}_{month}_{year}"
            if station_key in processed_stations:
                print(f"Skipping duplicate station: {station_id} for {month}/{year}")
                continue
            
            processed_stations[station_key] = True
            station_name = station_info["name"]
            
            # Calculate averages from sums and counts
            station_data = []
            for hour in hours:
                sum_val = station_info["sums"][hour]
                count = station_info["counts"][hour]
                if count > 0:
                    station_data.append({"AM_PM": hour, "ridership": sum_val / count})
            
            # Convert to DataFrame
            if not station_data:
                print(f"No data for station {station_id} in {month}/{year}")
                continue
                
            station_df = pd.DataFrame(station_data)
            
            # Skip stations with insufficient data
            if len(station_df) < 12:  # Require at least half the hours to have data
                print(f"Skipping station {station_id} due to insufficient data points ({len(station_df)} hours)")
                continue
                
            # Merge with all_hours to ensure all 24 hours are represented
            station_df = all_hours.merge(station_df, on="AM_PM", how="left")
            
            # Check if we have enough non-null data points after merging
            valid_data_points = station_df["ridership"].count()
            if valid_data_points < 12:
                print(f"Skipping station {station_id} due to insufficient valid data points ({valid_data_points} hours)")
                continue
                
            # Fill missing values
            station_df["ridership"] = station_df["ridership"].interpolate(method="linear")
            
            # Sort by hour for proper display
            station_df["hour_idx"] = station_df["AM_PM"].map(hour_indices)
            station_df = station_df.sort_values("hour_idx").drop("hour_idx", axis=1)
            
            # Sanitize station name
            sanitized_station_name = sanitize_name(station_name)
            
            try:
                # Create plot in memory with minimal memory usage
                plt.figure(figsize=(10, 6))
                plt.plot(station_df["AM_PM"], station_df["ridership"], marker='o', linestyle='-', 
                         label=f"Avg Riders: {station_df['ridership'].mean():.2f}")
                plt.title(f"Avg Ridership for {sanitized_station_name} - {month}/{year}", fontsize=16)
                plt.xlabel("Time (EST)", fontsize=14)
                plt.ylabel("Avg Riders", fontsize=14)
                plt.xticks(rotation=45)
                plt.grid(True, linestyle='--', alpha=0.6)
                plt.legend()
                
                # Add watermark with your name - positioned lower and to the right
                plt.figtext(0.45, -0.01, WATERMARK_TEXT, ha='center', color='gray', alpha=0.7, fontsize=10)
                
                plt.tight_layout()
                
                # Save to BytesIO instead of file
                img_bytes = io.BytesIO()
                plt.savefig(img_bytes, format='png', dpi=96)  # Lower DPI for memory savings
                img_bytes.seek(0)
                plt.close()  # Close to release memory
                
                # Add to PowerPoint
                slide = ppt.slides.add_slide(ppt.slide_layouts[5])
                title = slide.shapes.title
                title.text = f"{sanitized_station_name} - {month}/{year}"
                
                # Add the image directly from memory
                left = Inches(1)
                top = Inches(1.5)
                height = Inches(5)
                slide.shapes.add_picture(img_bytes, left, top, height=height)
                
                chart_count += 1
                
                # Clean up
                del img_bytes
                
            except Exception as e:
                print(f"Error creating chart for station {station_id}: {str(e)}")
                
            # Clear station data
            del station_df
            gc.collect()
        
        # Save PowerPoint to the specified output directory
        if chart_count > 0:
            ppt_filename = f"MTA_Ridership_{month}_{year}.pptx"
            ppt_path = os.path.join(file_path_OutPut, ppt_filename)
            ppt.save(ppt_path)
            print(f"PowerPoint generated with {chart_count} charts: {ppt_path}")
        else:
            print(f"No charts generated for {month}/{year}, skipping PowerPoint creation")
        
        # Clear memory
        del ppt
        gc.collect()
    
    # Clear all month data after processing the year
    del month_station_data
    gc.collect()

# Print summary
print(f"Total unique station-month combinations processed: {len(processed_stations)}")