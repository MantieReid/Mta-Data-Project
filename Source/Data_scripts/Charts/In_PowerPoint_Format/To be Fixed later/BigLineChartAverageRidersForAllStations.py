import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
import os

# Load the dataset
# Build path relative to current working directory
current_dir = os.getcwd()
# Navigate up from current directory to project root
project_root = os.path.abspath(os.path.join(current_dir, "..", "..", "..", ".."))
file_path = os.path.join(project_root, "Source", "Data", "Raw", "MTA_Subway_Hourly_Ridership__2020-2024.csv")
print(f"Looking for file at: {file_path}")

try:
    df = pd.read_csv(file_path)
except FileNotFoundError:
    print(f"File not found: {file_path}")
    print(f"Current working directory: {os.getcwd()}")
    # List files in the expected directory to debug
    expected_dir = os.path.dirname(file_path)
    if os.path.exists(expected_dir):
        print(f"Files in {expected_dir}:")
        print(os.listdir(expected_dir))
    else:
        print(f"Directory does not exist: {expected_dir}")
    exit(1)

# Convert transit_timestamp to datetime format
df["transit_timestamp"] = pd.to_datetime(df["transit_timestamp"], errors="coerce")

# Extract year and hour
df["year"] = df["transit_timestamp"].dt.year
df["hour"] = df["transit_timestamp"].dt.hour

# Convert military time to 12-hour AM/PM format
df["AM_PM"] = df["hour"].apply(lambda x: f"{x % 12 if x % 12 != 0 else 12} {'AM' if x < 12 else 'PM'}")

# Define all possible hours for completeness
all_hours = pd.DataFrame({"AM_PM": [f"{h % 12 if h % 12 != 0 else 12} {'AM' if h < 12 else 'PM'}" for h in range(24)]})

# Filter data for 2023 and 2024
df_filtered = df[df["year"].isin([2023, 2024])]

# Calculate average ridership per hour across all stations
avg_ridership = df_filtered.groupby(["year", "AM_PM"])['ridership'].mean().reset_index()

# Merge with all_hours to ensure all 24 hours are represented
avg_ridership = all_hours.merge(avg_ridership, on="AM_PM", how="left")

# Pivot for plotting
avg_ridership_pivot = avg_ridership.pivot(index="AM_PM", columns="year", values="ridership")

# Create line chart
plt.figure(figsize=(10, 6))
for year in [2023, 2024]:
    plt.plot(avg_ridership_pivot.index, avg_ridership_pivot[year], marker='o', linestyle='-', label=f"{year}")

plt.title("Average Ridership Per Hour (All Stations) - 2023 vs 2024", fontsize=16)
plt.xlabel("Time (EST)", fontsize=14)
plt.ylabel("Avg Riders", fontsize=14)
plt.xticks(rotation=45)
plt.grid(True, linestyle='--', alpha=0.6)
plt.legend()
plt.tight_layout()

# Save chart as an image
chart_filename = "Average_Ridership_Per_Hour_2023_2024.png"
plt.savefig(chart_filename)
plt.close()

# Create PowerPoint presentation
ppt = Presentation()
slide = ppt.slides.add_slide(ppt.slide_layouts[5])  # Blank slide

# Add title
title = slide.shapes.title
if title:
    title.text = "Average Ridership Per Hour (2023 vs 2024)"

# Add chart image to PowerPoint
left = Inches(1)
top = Inches(1.5)
height = Inches(5)
slide.shapes.add_picture(chart_filename, left, top, height=height)

# Save PowerPoint file
ppt_path = "Average_Ridership_Presentation.pptx"
ppt.save(ppt_path)

print(f"PowerPoint generated: {ppt_path}")